# -*- coding: utf-8 -*-
"""PIE_가이드.html 의 3개 국어 내용을 그대로 읽어 PowerPoint 파일 3개를 만든다.
   가이드 HTML 이 유일한 원본이고, 이 스크립트는 옮기기만 한다."""
import io, re, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SRC = r'D:\김형래\프로그램\작업분석 중국\PIE_가이드.html'
OUT = r'D:\김형래\프로그램\작업분석 중국'

# ── 1) HTML 안의 var G = {...}; 를 JSON 으로 변환 ────────────────────────────
def js_object_to_json(js):
    out, i, n = [], 0, len(js)
    while i < n:
        c = js[i]
        if c == "'":                      # 작은따옴표 문자열 → JSON 문자열
            i += 1; buf = []
            while i < n and js[i] != "'":
                if js[i] == '\\':
                    buf.append(js[i]); i += 1
                    if i < n: buf.append(js[i]); i += 1
                    continue
                buf.append(js[i]); i += 1
            i += 1
            s = ''.join(buf).replace('"', '\\"')
            out.append('"' + s + '"')
        elif c == '/' and i + 1 < n and js[i+1] == '*':   # 주석 제거
            j = js.find('*/', i + 2); i = (j + 2) if j >= 0 else n
        elif c == '/' and i + 1 < n and js[i+1] == '/':
            j = js.find('\n', i); i = (j + 1) if j >= 0 else n
        else:
            out.append(c); i += 1
    t = ''.join(out)
    t = re.sub(r'([{,\[]\s*)([A-Za-z_][A-Za-z0-9_]*)\s*:', r'\1"\2":', t)  # 키에 따옴표
    t = re.sub(r',(\s*[}\]])', r'\1', t)                                    # 마지막 쉼표 제거
    return t

src = io.open(SRC, encoding='utf-8').read()
a = src.index('var G = {')
b = src.index('\n};', a)
import json
G = json.loads(js_object_to_json(src[a + len('var G = '): b + 2]))

# ── 2) 색·글꼴 ──────────────────────────────────────────────────────────────
BG     = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x11, 0x18, 0x27)
SUB    = RGBColor(0x47, 0x55, 0x69)
DIM    = RGBColor(0x94, 0xA3, 0xB8)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
GREEN  = RGBColor(0x15, 0x9E, 0x4E)
AMBER  = RGBColor(0xB4, 0x53, 0x09)
RED    = RGBColor(0xDC, 0x26, 0x26)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
PANEL  = RGBColor(0xF1, 0xF5, 0xF9)
BORDER = RGBColor(0xCB, 0xD5, 0xE1)
DARK   = RGBColor(0x0F, 0x17, 0x2A)   # 화면 모형용
DARK2  = RGBColor(0x1E, 0x29, 0x3B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT = {'ko': '맑은 고딕', 'zh': '微软雅黑', 'vi': 'Segoe UI'}
FNAME = {'ko': 'PIE_가이드_한국어.pptx', 'zh': 'PIE_가이드_中文.pptx', 'vi': 'PIE_가이드_TiengViet.pptx'}

# ── 3) HTML 조각 → (텍스트, 굵게, 색) 런 목록 ───────────────────────────────
TAG = re.compile(r'<(/?)(b|code|span)([^>]*)>')
def runs_of(html):
    res, pos, bold, col = [], 0, 0, []
    for m in TAG.finditer(html):
        if m.start() > pos:
            res.append((html[pos:m.start()], bold > 0, col[-1] if col else None))
        closing, tag, attr = m.group(1), m.group(2), m.group(3)
        if tag in ('b', 'code'):
            bold += -1 if closing else 1
            if tag == 'code' and not closing: col.append(BLUE)
            elif tag == 'code' and closing and col: col.pop()
        elif tag == 'span':
            if closing:
                if col: col.pop()
            else:
                c = re.search(r'#([0-9a-fA-F]{6})', attr)
                col.append(RGBColor.from_string(c.group(1).upper()) if c else INK)
        pos = m.end()
    if pos < len(html):
        res.append((html[pos:], bold > 0, col[-1] if col else None))
    return [(re.sub(r'<[^>]+>', '', t), bd, c) for t, bd, c in res if t]

def put(tf, html, size=13, color=INK, font='맑은 고딕', space_after=4, bullet_col=None):
    p = tf.add_paragraph() if tf.paragraphs[0].runs or tf.paragraphs[0].text else tf.paragraphs[0]
    p.space_after = Pt(space_after)
    for t, bd, c in runs_of(html):
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = bd
        r.font.color.rgb = c if c else color
        r.font.name = font
    return p

def box(sl, x, y, w, h, fill=None, line=None, rad=False):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rad else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    if rad:
        try: sh.adjustments[0] = 0.08
        except Exception: pass
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh

def label(sl, x, y, w, h, text, size=12, color=INK, bold=False, font='맑은 고딕',
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = align
    for t, bd, c in runs_of(text):
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = bold or bd
        r.font.color.rgb = c if c else color; r.font.name = font
    return tb

SW, SH = 13.333, 7.5

def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
    return prs

def blank(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill; bg.solid(); bg.fore_color.rgb = BG
    return sl

def head(sl, title, sub, font, page=None):
    box(sl, 0, 0, SW, 0.92, fill=RGBColor(0x0F, 0x23, 0x3F))
    label(sl, 0.55, 0.14, 10.5, 0.4, title, size=22, color=WHITE, bold=True, font=font)
    if sub: label(sl, 0.55, 0.55, 11.5, 0.3, sub, size=11.5, color=RGBColor(0x93,0xC5,0xFD), font=font)
    if page: label(sl, SW-1.6, 0.28, 1.1, 0.35, page, size=12, color=RGBColor(0x93,0xC5,0xFD),
                   font=font, align=PP_ALIGN.RIGHT)

# ── 4) 네이티브 도해 ────────────────────────────────────────────────────────
def fig_install(sl, F, font, y0=1.4):
    steps = [('📁', F['f_copy']), ('🖱', F['f_dbl']), ('⬛', F['f_black']), ('🌐', F['f_browser'])]
    x = 0.75
    for i, (ic, tx) in enumerate(steps):
        box(sl, x, y0, 2.7, 1.5, fill=PANEL, line=BORDER, rad=True)
        label(sl, x, y0+0.18, 2.7, 0.5, ic, size=26, font=font, align=PP_ALIGN.CENTER)
        label(sl, x+0.15, y0+0.78, 2.4, 0.55, tx, size=12.5, bold=True, font=font, align=PP_ALIGN.CENTER)
        if i < 3:
            label(sl, x+2.75, y0+0.55, 0.5, 0.4, '▶', size=20, color=BLUE, font=font, align=PP_ALIGN.CENTER)
        x += 3.13
    label(sl, 0.75, y0+1.7, 11.8, 0.35, F['f_blackWarn'], size=12.5, color=RED, bold=True,
          font=font, align=PP_ALIGN.CENTER)

def fig_marking(sl, F, font, y0=1.3):
    # 왼쪽: 영상 + 클릭 두 번
    box(sl, 0.6, y0, 5.6, 3.0, fill=DARK, line=DARK2, rad=True)
    label(sl, 0.6, y0+1.15, 5.6, 0.5, '▶', size=30, color=RGBColor(0x1E,0x3A,0x5F), font=font, align=PP_ALIGN.CENTER)
    label(sl, 0.9, y0+1.9, 5.0, 0.35, F['f_videoHere'], size=12, color=DIM, font=font, align=PP_ALIGN.CENTER)
    for cx, cy, n, col, tx in ((2.2, y0+0.95, '1', GREEN, F['f_clickStart']),
                               (4.6, y0+2.05, '2', RED, F['f_clickEnd'])):
        d = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(0.34), Inches(0.34))
        d.fill.solid(); d.fill.fore_color.rgb = col; d.line.fill.background(); d.shadow.inherit = False
        label(sl, cx, cy, 0.34, 0.34, n, size=12, color=WHITE, bold=True, font=font, align=PP_ALIGN.CENTER)
        label(sl, cx-1.0, cy-0.42, 2.4, 0.3, tx, size=11.5, color=col, bold=True, font=font, align=PP_ALIGN.CENTER)
    # 손 버튼
    label(sl, 0.6, y0+3.15, 0.6, 0.3, F['f_hand'] + ':', size=12, color=SUB, font=font)
    hx = 1.25
    for i, (t, c) in enumerate(((F['f_lh'], BLUE), (F['f_rh'], GREEN), (F['f_one'], AMBER))):
        box(sl, hx, y0+3.12, 1.62, 0.36, fill=(c if i == 0 else None), line=c, rad=True)
        label(sl, hx, y0+3.12, 1.62, 0.36, t, size=11, color=(WHITE if i == 0 else c), bold=True,
              font=font, align=PP_ALIGN.CENTER)
        hx += 1.72
    # 오른쪽: 표 (활성 손 + 구분선 + 다른 손)
    tx0, tw = 6.55, 6.2
    box(sl, tx0, y0, tw, 3.55, fill=DARK, line=DARK2, rad=True)
    box(sl, tx0+0.1, y0+0.1, tw-0.2, 0.32, fill=RGBColor(0x16,0x23,0x3A))
    cols = [(0.25, F['f_taskName']), (2.1, F['f_partCol']), (3.7, F['f_type']), (4.9, F['f_startC'])]
    for dx, t in cols:
        label(sl, tx0+dx, y0+0.1, 1.8, 0.32, t, size=10, color=RGBColor(0x93,0xC5,0xFD), bold=True, font=font)
    names = [F['f_egPick'], F['f_egMove'], F['f_egFit']]
    for r in range(3):
        yy = y0 + 0.5 + r*0.38
        label(sl, tx0+0.25, yy, 1.8, 0.34, names[r], size=11, color=WHITE, font=font)
        label(sl, tx0+2.1, yy, 1.5, 0.34, 'A-100%d' % (r+1), size=10.5, color=RGBColor(0x86,0xEF,0xAC), font=font)
        label(sl, tx0+3.7, yy, 1.1, 0.34, F['f_va'], size=10.5, color=RGBColor(0x4A,0xDE,0x80), font=font)
        label(sl, tx0+4.9, yy, 1.1, 0.34, '00:0%d.2' % (r+1), size=10, color=DIM, font=font)
    dv = y0 + 1.72
    box(sl, tx0+0.08, dv, tw-0.16, 0.03, fill=GREEN)
    label(sl, tx0+0.25, dv+0.05, 1.6, 0.3, F['f_rh'], size=11, color=RGBColor(0x4A,0xDE,0x80), bold=True, font=font)
    label(sl, tx0+1.9, dv+0.05, 4.1, 0.3, F['f_otherHand'], size=10, color=DIM, font=font)
    for r in range(2):
        yy = dv + 0.42 + r*0.38
        box(sl, tx0+0.08, yy-0.02, tw-0.16, 0.36, fill=RGBColor(0x0B,0x12,0x20))
        label(sl, tx0+0.25, yy, 1.8, 0.34, [F['f_egPick'], F['f_egInsert']][r], size=11,
              color=RGBColor(0x64,0x74,0x8B), font=font)
        label(sl, tx0+2.1, yy, 1.5, 0.34, 'A-100%d' % (r+4), size=10.5, color=RGBColor(0x3F,0x6B,0x52), font=font)

def fig_pitch(sl, F, font, y0=1.5):
    vals = [160, 124, 118, 104]
    x0, w, axis = 2.2, 8.2, 160.0
    tak = x0 + w*100/axis
    label(sl, 0.6, y0-0.45, 6.0, 0.35, '📏 ' + F['f_pitch'], size=15, bold=True, font=font)
    for i, v in enumerate(vals):
        y = y0 + i*0.72
        label(sl, 0.6, y, 1.5, 0.4, F['f_worker'] + ' %d' % (i+1), size=12, color=SUB, font=font)
        box(sl, x0, y+0.06, w, 0.32, fill=RGBColor(0xE2,0xE8,0xF0), rad=True)
        col = RED if v > 100 else (AMBER if v > 90 else GREEN)
        box(sl, x0, y+0.06, w*min(v,100)/axis, 0.32, fill=col, rad=True)
        if v > 100:
            ov = box(sl, tak, y+0.06, w*(v-100)/axis, 0.32, fill=RGBColor(0x99,0x1B,0x1B), rad=True)
            ov.line.color.rgb = RED; ov.line.width = Pt(1)
        label(sl, x0+w+0.15, y+0.06, 0.9, 0.32, '%d%%' % v, size=12.5, color=col, bold=True, font=font)
    ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(tak), Inches(y0-0.05),
                             Inches(0.015), Inches(len(vals)*0.72+0.05))
    ln.fill.solid(); ln.fill.fore_color.rgb = INK; ln.line.fill.background(); ln.shadow.inherit = False
    label(sl, tak-0.5, y0-0.42, 1.0, 0.3, '100%', size=11, bold=True, font=font, align=PP_ALIGN.CENTER)
    label(sl, 0.6, y0+len(vals)*0.72+0.15, 12.1, 0.4,
          '① ' + F['f_cTakt'] + '    ② ' + F['f_cPitch'], size=12.5, color=AMBER, bold=True, font=font)

def fig_hand(sl, F, font, y0=1.5):
    rows = [(F['f_lh'], '3.4s', '34%', BLUE, 3.4),
            (F['f_rh'], '2.5s', '25%', GREEN, 2.5),
            (F['f_one'], '4.1s', '41%', AMBER, 4.1)]
    label(sl, 0.6, y0-0.5, 6.0, 0.35, '⚖ ' + F['f_handEval'], size=15, bold=True, font=font)
    label(sl, 4.2, y0-0.5, 2.0, 0.35, F['f_hbr'], size=12, color=SUB, font=font)
    label(sl, 5.6, y0-0.55, 1.2, 0.45, '74%', size=20, color=AMBER, bold=True, font=font)
    box(sl, 7.0, y0-0.5, 5.6, 0.38, fill=RGBColor(0xFE,0xE2,0xE2), line=RED, rad=True)
    label(sl, 7.15, y0-0.5, 5.3, 0.38, '⚠ ' + F['f_hbrWarn'], size=11.5, color=RED, font=font)
    hdr = [(0.6, F['f_hand']), (2.5, F['f_total']), (3.5, F['f_share']), (4.5, F['f_va']),
           (5.5, F['f_aux']), (6.5, F['f_waste'])]
    for dx, t in hdr:
        label(sl, dx, y0, 1.6, 0.32, t, size=11, color=BLUE, bold=True, font=font)
    for i, (nm, tot, sh, col, val) in enumerate(rows):
        y = y0 + 0.42 + i*0.5
        label(sl, 0.6, y, 1.9, 0.36, nm, size=12, color=col, bold=True, font=font)
        label(sl, 2.5, y, 1.0, 0.36, tot, size=12, bold=True, font=font)
        label(sl, 3.5, y, 1.0, 0.36, sh, size=12, color=SUB, font=font)
        label(sl, 4.5, y, 1.0, 0.36, tot, size=12, color=GREEN, font=font)
        label(sl, 5.5, y, 1.0, 0.36, '0.0s', size=12, color=AMBER, font=font)
        label(sl, 6.5, y, 1.0, 0.36, '0.0s', size=12, color=DIM, font=font)
        box(sl, 7.6, y+0.08, 5.0, 0.2, fill=RGBColor(0xE2,0xE8,0xF0), rad=True)
        box(sl, 7.6, y+0.08, 5.0*val/4.1, 0.2, fill=col, rad=True)
    label(sl, 0.6, y0+0.42+len(rows)*0.5+0.15, 12.1, 0.4, '① ' + F['f_cHandEval'],
          size=12.5, color=AMBER, bold=True, font=font)

def _screen(sl, F, font, x, y, w, h, active):
    """왼쪽 메뉴가 있는 앱 화면 모형"""
    box(sl, x, y, w, h, fill=DARK, line=DARK2, rad=True)
    box(sl, x+0.06, y+0.06, 1.85, h-0.12, fill=RGBColor(0x0B,0x12,0x20))
    for i, m in enumerate(F['menu']):
        yy = y + 0.16 + i*0.31
        if yy + 0.28 > y + h: break
        if i == active:
            box(sl, x+0.06, yy-0.02, 1.85, 0.30, fill=RGBColor(0x12,0x23,0x3F))
            box(sl, x+0.06, yy-0.02, 0.05, 0.30, fill=BLUE)
        label(sl, x+0.22, yy, 1.6, 0.28, m, size=10,
              color=(WHITE if i == active else RGBColor(0x64,0x74,0x8B)),
              bold=(i == active), font=font)
    return x+1.98

def fig_home(sl, F, font, y0=1.3):
    x, w, h = 0.9, 11.5, 4.0
    cx = _screen(sl, F, font, x, y0, w, h, 0)
    label(sl, cx, y0+0.5, w-2.2, 0.6, 'PIE', size=30, color=RGBColor(0x60,0xA5,0xFA), bold=True,
          font=font, align=PP_ALIGN.CENTER)
    label(sl, cx, y0+1.1, w-2.2, 0.4, 'Powernet Industrial Engineering', size=12,
          color=RGBColor(0x94,0xA3,0xB8), font=font, align=PP_ALIGN.CENTER)
    bw = 5.4; bx = cx + (w-2.2-bw)/2
    box(sl, bx, y0+1.7, bw, 0.75, fill=BLUE, rad=True)
    label(sl, bx, y0+1.7, bw, 0.75, '📹 ' + F['f_loadVideo'], size=17, color=WHITE, bold=True,
          font=font, align=PP_ALIGN.CENTER)
    label(sl, bx, y0+2.7, bw, 0.5, F['f_noSession'], size=11.5, color=RGBColor(0x64,0x74,0x8B),
          font=font, align=PP_ALIGN.CENTER)
    label(sl, 0.9, y0+h+0.12, 12.1, 0.4,
          '① ' + F['f_cMenu'] + '    ② ' + F['f_cLoad'] + '    ③ ' + F['f_cRating'],
          size=12.5, color=AMBER, bold=True, font=font)

def fig_parts(sl, F, font, y0=1.3):
    x, w, h = 0.9, 11.5, 3.9
    cx = _screen(sl, F, font, x, y0, w, h, 8)
    bx = cx + 0.1
    for t, c in ((F['f_bImport'], BLUE), (F['f_bClean'], AMBER), (F['f_bReset'], RGBColor(0x7F,0x1D,0x1D))):
        box(sl, bx, y0+0.14, 2.35, 0.34, fill=c, rad=True)
        label(sl, bx, y0+0.14, 2.35, 0.34, t, size=10.5, color=WHITE, bold=True, font=font,
              align=PP_ALIGN.CENTER)
        bx += 2.5
    box(sl, cx+0.1, y0+0.62, 2.9, h-0.78, fill=RGBColor(0x0B,0x12,0x20), line=DARK2)
    for k in range(4):
        yy = y0 + 0.74 + k*0.72
        box(sl, cx+0.2, yy, 2.7, 0.62, fill=(RGBColor(0x12,0x23,0x3F) if k == 0 else None),
            line=(BLUE if k == 0 else DARK2), rad=True)
        label(sl, cx+0.32, yy+0.03, 2.5, 0.28, 'A-100%d(CAP-ELEC)' % (k+1), size=10.5,
              color=(WHITE if k == 0 else RGBColor(0x94,0xA3,0xB8)), bold=True, font=font)
        label(sl, cx+0.32, yy+0.3, 2.5, 0.26, '400V,10uF,105C', size=9,
              color=RGBColor(0x64,0x74,0x8B), font=font)
    dx = cx + 3.15
    box(sl, dx, y0+0.62, w-(dx-x)-0.15, h-0.78, fill=RGBColor(0x0B,0x12,0x20), line=DARK2)
    label(sl, dx+0.15, y0+0.75, 4.0, 0.32, 'A-1001(CAP-ELEC)', size=13, color=WHITE, bold=True, font=font)
    box(sl, dx+0.12, y0+1.14, w-(dx-x)-0.4, 0.5, fill=DARK, line=DARK2, rad=True)
    px = dx+0.25
    for t, v in ((F['f_code'], 'A-1001'), (F['f_pname'], 'CAP-ELEC'), (F['f_spec'], '400V,10uF')):
        label(sl, px, y0+1.14, 0.8, 0.5, t, size=9.5, color=RGBColor(0x64,0x74,0x8B), font=font)
        label(sl, px+0.72, y0+1.14, 1.6, 0.5, v, size=10, color=WHITE, bold=True, font=font)
        px += 2.3
    label(sl, 0.9, y0+h+0.12, 12.1, 0.4,
          '① ' + F['f_cImport'] + '    ② ' + F['f_cList'] + '    ③ ' + F['f_cBom'],
          size=12.5, color=AMBER, bold=True, font=font)

def fig_settings(sl, F, font, y0=1.4):
    box(sl, 3.4, y0, 6.5, 2.55, fill=PANEL, line=BORDER, rad=True)
    label(sl, 3.7, y0+0.15, 5.0, 0.4, '⚙️ ' + F['f_settings'], size=15, bold=True, font=font)
    box(sl, 3.65, y0+0.6, 6.0, 1.55, fill=RGBColor(0xE8,0xF1,0xFE), line=BLUE, rad=True)
    label(sl, 3.85, y0+0.7, 5.6, 0.35, '⏱ ' + F['f_stBasis'], size=13, color=BLUE, bold=True, font=font)
    label(sl, 3.85, y0+1.05, 5.6, 0.3, F['f_formula'], size=11, color=SUB, font=font)
    label(sl, 3.85, y0+1.45, 1.0, 0.34, F['ratingL'], size=12, font=font)
    box(sl, 4.95, y0+1.45, 0.85, 0.34, fill=WHITE, line=BORDER)
    label(sl, 4.95, y0+1.45, 0.85, 0.34, '100 %', size=12, bold=True, font=font, align=PP_ALIGN.CENTER)
    label(sl, 6.1, y0+1.45, 1.1, 0.34, F['allowL'], size=12, font=font)
    box(sl, 7.3, y0+1.45, 0.85, 0.34, fill=WHITE, line=BORDER)
    label(sl, 7.3, y0+1.45, 0.85, 0.34, '15 %', size=12, bold=True, font=font, align=PP_ALIGN.CENTER)
    label(sl, 3.85, y0+1.85, 5.6, 0.3, F['f_egCalc'], size=12, color=BLUE, bold=True, font=font)
    label(sl, 3.4, y0+2.7, 6.5, 0.4, '① ' + F['f_cFormula'] + '    ② ' + F['f_cInput'],
          size=12.5, color=AMBER, bold=True, font=font, align=PP_ALIGN.CENTER)

def fig_chart(sl, F, font, y0=1.4):
    base = y0 + 2.9
    segs = ((BLUE, 0.52, F['f_egPick']), (VIOLET, 0.30, F['f_egMove']), (RGBColor(0xDB,0x27,0x77), 0.86, F['f_egFit']))
    xs, bw = (2.0, 5.3, 8.6), 1.7
    tops = []
    for ci, x in enumerate(xs):
        y, arr = base, []
        for col, hh0, nm in segs:
            hh = hh0 * (1 + ci*0.12); y -= hh
            box(sl, x, y, bw, hh-0.03, fill=col)
            if hh > 0.34:
                label(sl, x, y, bw, hh, nm, size=11, color=WHITE, bold=True, font=font, align=PP_ALIGN.CENTER)
            arr.append((y, hh))
        tops.append(arr)
        label(sl, x, base+0.08, bw, 0.35, 'C%d' % (ci+1), size=13,
              color=(BLUE, GREEN, AMBER)[ci], bold=True, font=font, align=PP_ALIGN.CENTER)
        label(sl, x, y-0.35, bw, 0.3, '%.1fs' % (31.6+ci*1.4), size=12, bold=True, font=font,
              align=PP_ALIGN.CENTER)
    for g in range(2):
        for si in range(3):
            (y1, h1), (y2, h2) = tops[g][si], tops[g+1][si]
            x1, x2 = xs[g]+bw, xs[g+1]
            ff = sl.shapes.build_freeform(Inches(x1), Inches(y1))
            ff.add_line_segments([(Inches(x2), Inches(y2)), (Inches(x2), Inches(y2+h2)),
                                  (Inches(x1), Inches(y1+h1))], close=True)
            sh = ff.convert_to_shape()
            sh.fill.solid(); sh.fill.fore_color.rgb = segs[si][0]
            sh.fill.transparency = 0.82
            sh.line.fill.background(); sh.shadow.inherit = False
    label(sl, 0.7, base+0.5, 12.1, 0.4, '① ' + F['f_cBandDesc'], size=12.5, color=AMBER,
          bold=True, font=font)

def fig_line(sl, F, font, y0=1.3):
    box(sl, 0.7, y0, 5.4, 2.5, fill=DARK, line=DARK2, rad=True)
    label(sl, 0.7, y0+0.9, 5.4, 0.5, '🎬', size=26, font=font, align=PP_ALIGN.CENTER)
    label(sl, 0.7, y0+1.5, 5.4, 0.35, F['f_pickWorkerVideo'], size=12, color=DIM, font=font,
          align=PP_ALIGN.CENTER)
    bx = 0.7
    for t, c in ((F['f_markStart'], GREEN), (F['f_markEnd'], RED), (F['f_nextCycle'] + ' +', BLUE)):
        box(sl, bx, y0+2.65, 1.72, 0.4, fill=None, line=c, rad=True)
        label(sl, bx, y0+2.65, 1.72, 0.4, t, size=11, color=c, bold=True, font=font, align=PP_ALIGN.CENTER)
        bx += 1.85
    tx = 6.4
    box(sl, tx, y0, 6.25, 3.05, fill=DARK, line=DARK2, rad=True)
    label(sl, tx+0.2, y0+0.1, 4.0, 0.35, F['f_markList'] + '  ' + F['f_lh'] + ' (4)', size=12,
          color=RGBColor(0x60,0xA5,0xFA), bold=True, font=font)
    for r in range(3):
        yy = y0 + 0.55 + r*0.42
        label(sl, tx+0.2, yy, 0.3, 0.34, str(r+1), size=10.5, color=DIM, font=font)
        box(sl, tx+0.55, yy, 3.0, 0.34, fill=RGBColor(0x0B,0x12,0x20), line=RGBColor(0x14,0x53,0x2D))
        label(sl, tx+0.65, yy, 2.9, 0.34, 'A-100%d(CAP-ELEC) ▾' % (r+1), size=10.5,
              color=RGBColor(0x86,0xEF,0xAC), font=font)
        label(sl, tx+3.75, yy, 0.6, 0.34, 'C1', size=10.5, color=RGBColor(0x60,0xA5,0xFA), bold=True, font=font)
        label(sl, tx+4.45, yy, 0.6, 0.34, 'LH', size=10.5, color=RGBColor(0x60,0xA5,0xFA), font=font)
    dv = y0 + 1.9
    box(sl, tx+0.12, dv, 6.0, 0.03, fill=GREEN)
    label(sl, tx+0.2, dv+0.06, 1.6, 0.3, F['f_rh'], size=11, color=RGBColor(0x4A,0xDE,0x80),
          bold=True, font=font)
    label(sl, tx+1.9, dv+0.06, 4.2, 0.3, F['f_otherHand'], size=10, color=DIM, font=font)
    for r in range(2):
        yy = dv + 0.44 + r*0.4
        box(sl, tx+0.12, yy-0.02, 6.0, 0.38, fill=RGBColor(0x0B,0x12,0x20))
        label(sl, tx+0.65, yy, 2.9, 0.34, 'A-100%d(FUSE) ▾' % (r+4), size=10.5,
              color=RGBColor(0x3F,0x6B,0x52), font=font)
    label(sl, 0.7, y0+3.25, 12.1, 0.4,
          '① ' + F['f_cPickPart'] + '    ② ' + F['f_cInherit'], size=12.5, color=AMBER,
          bold=True, font=font)

def fig_sim(sl, F, font, y0=1.4):
    k = ((F['f_ct'], '16.7s', RED), (F['f_bottleneck'], F['f_worker'] + ' 1', AMBER),
         (F['f_balRate'], '79%', BLUE), (F['f_perHour'], '215' + F['f_ea'], GREEN))
    for i, (t, v, c) in enumerate(k):
        x = 0.7 + i*3.1
        box(sl, x, y0, 2.9, 0.85, fill=PANEL, line=BORDER, rad=True)
        label(sl, x+0.2, y0+0.08, 2.5, 0.3, t, size=11, color=SUB, font=font)
        label(sl, x+0.2, y0+0.38, 2.5, 0.42, v, size=18, color=c, bold=True, font=font)
    base = y0 + 3.3
    for i, (hh, hot) in enumerate(((1.55, 1), (1.0, 0), (0.95, 0), (0.82, 0))):
        x = 1.6 + i*2.6
        if not hot:
            box(sl, x, base-1.55, 1.9, 1.55-hh, fill=RGBColor(0xE2,0xE8,0xF0))
            label(sl, x, base-1.55, 1.9, 1.55-hh, F['f_idle'], size=11, color=SUB, font=font,
                  align=PP_ALIGN.CENTER)
        b = box(sl, x, base-hh, 1.9, hh, fill=(RED if hot else GREEN))
        if hot:
            b.line.color.rgb = RGBColor(0x7F,0x1D,0x1D); b.line.width = Pt(2)
            label(sl, x-0.3, base-hh-0.38, 2.5, 0.32, '⚠ ' + F['f_bottleneck'], size=12,
                  color=RED, bold=True, font=font, align=PP_ALIGN.CENTER)
        label(sl, x, base+0.08, 1.9, 0.32, F['f_worker'] + ' %d' % (i+1), size=11.5, color=SUB,
              bold=True, font=font, align=PP_ALIGN.CENTER)
    label(sl, 0.7, base+0.55, 12.1, 0.4, '※ ' + F['f_simNote'], size=12.5, color=AMBER,
          bold=True, font=font)

def fig_predict(sl, F, font, y0=1.4):
    box(sl, 0.7, y0, 4.3, 2.9, fill=PANEL, line=BORDER, rad=True)
    label(sl, 0.95, y0+0.12, 3.8, 0.35, F['f_bomInput'], size=13, color=BLUE, bold=True, font=font)
    box(sl, 0.95, y0+0.55, 3.8, 0.38, fill=WHITE, line=BORDER)
    label(sl, 1.1, y0+0.55, 3.6, 0.38, F['f_pickPart'] + ' ▾', size=11.5, color=SUB, font=font)
    label(sl, 0.95, y0+1.05, 0.9, 0.35, F['f_qty'], size=11.5, color=SUB, font=font)
    box(sl, 1.8, y0+1.05, 0.7, 0.35, fill=WHITE, line=BORDER)
    label(sl, 1.8, y0+1.05, 0.7, 0.35, '2', size=12, bold=True, font=font, align=PP_ALIGN.CENTER)
    box(sl, 2.65, y0+1.05, 1.0, 0.35, fill=BLUE, rad=True)
    label(sl, 2.65, y0+1.05, 1.0, 0.35, '+ ' + F['f_add'], size=11.5, color=WHITE, bold=True,
          font=font, align=PP_ALIGN.CENTER)
    for r in range(3):
        yy = y0 + 1.6 + r*0.42
        box(sl, 0.95, yy, 3.8, 0.36, fill=WHITE, line=BORDER)
        label(sl, 1.1, yy, 2.4, 0.36, 'A-100%d(CAP-ELEC)' % (r+1), size=11, font=font)
        label(sl, 3.5, yy, 0.5, 0.36, '×%d' % (r+1), size=11, color=SUB, font=font)
        label(sl, 4.0, yy, 0.7, 0.36, '%.1fs' % (4.2+r), size=11, color=VIOLET, bold=True, font=font)
    ox = 5.35
    for i, (t, v, c) in enumerate(((F['f_expCt'], '22.4s', BLUE), (F['f_needPeople'], '5' + F['f_person'], AMBER),
                                   (F['f_balRate'], '86%', GREEN))):
        x = ox + i*2.5
        box(sl, x, y0, 2.3, 0.95, fill=PANEL, line=BORDER, rad=True)
        label(sl, x+0.18, y0+0.1, 2.0, 0.3, t, size=11, color=SUB, font=font)
        label(sl, x+0.18, y0+0.42, 2.0, 0.45, v, size=19, color=c, bold=True, font=font)
    label(sl, ox, y0+1.15, 5.0, 0.35, F['f_procPlan'], size=12.5, color=BLUE, bold=True, font=font)
    for p in range(3):
        yy = y0 + 1.55 + p*0.45
        box(sl, ox, yy, 7.3, 0.38, fill=WHITE, line=BORDER)
        label(sl, ox+0.15, yy, 1.3, 0.38, F['f_proc'] + ' %d' % (p+1), size=11, color=SUB, bold=True, font=font)
        box(sl, ox+1.5, yy+0.05, 2.0+p*0.6, 0.28, fill=GREEN, rad=True)
        label(sl, ox+1.65, yy, 3.5, 0.38, 'A-100%d → A-100%d' % (p+1, p+2), size=10, color=WHITE, font=font)

def fig_trouble(sl, F, font, y0=1.6):
    box(sl, 0.7, y0, 11.9, 1.1, fill=RGBColor(0x7F,0x1D,0x1D), rad=True)
    label(sl, 0.95, y0+0.1, 0.5, 0.4, '⚠', size=20, color=RGBColor(0xFE,0xE2,0xE2), font=font)
    label(sl, 1.45, y0+0.12, 8.5, 0.4, F['f_warnTitle'], size=13.5, color=RGBColor(0xFE,0xE2,0xE2),
          bold=True, font=font)
    label(sl, 1.45, y0+0.58, 8.5, 0.4, F['f_warnBody'], size=11.5, color=RGBColor(0xFE,0xCA,0xCA), font=font)
    box(sl, 10.2, y0+0.3, 2.1, 0.5, fill=BLUE, rad=True)
    label(sl, 10.2, y0+0.3, 2.1, 0.5, F['f_btnStdAddr'], size=11.5, color=WHITE, bold=True,
          font=font, align=PP_ALIGN.CENTER)
    label(sl, 0.7, y0+1.35, 12.1, 0.4, '※ ' + F['f_warnNote'], size=12.5, color=AMBER, bold=True, font=font)

FIGS = {'install': fig_install, 'home': fig_home, 'parts': fig_parts,
        'analysis': fig_marking, 'settings': fig_settings, 'chart': fig_chart, 'stats': fig_hand,
        'line': fig_line, 'balance': fig_pitch, 'sim': fig_sim, 'predict': fig_predict,
        'trouble': fig_trouble}
# 'parts_detail' 은 parts 그림 안에 상세 패널이 이미 들어 있어 따로 그리지 않는다
FIG_H = {'analysis': 4.1, 'home': 4.7, 'parts': 4.6, 'line': 4.1, 'chart': 3.9, 'predict': 3.5,
         'sim': 4.3, 'balance': 3.7, 'stats': 3.1, 'install': 2.5, 'settings': 3.4, 'trouble': 2.3}

# ── 5) 본문 슬라이드 ────────────────────────────────────────────────────────
def add_table(sl, hd, rows, x, y, w, font, fs=11):
    r, c = len(rows)+1, len(hd)
    h = 0.34 + 0.32*len(rows)
    gt = sl.shapes.add_table(r, c, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for j, t in enumerate(hd):
        cell = gt.cell(0, j); cell.text = ''
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = re.sub(r'<[^>]+>', '', t)
        run.font.size = Pt(fs); run.font.bold = True; run.font.name = font
        run.font.color.rgb = WHITE
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
        cell.margin_left = cell.margin_right = Inches(0.06)
    for i, row in enumerate(rows):
        for j, t in enumerate(row):
            cell = gt.cell(i+1, j); cell.text = ''
            p = cell.text_frame.paragraphs[0]; p.line_spacing = 0.95
            for tt, bd, cc in runs_of(t):
                run = p.add_run(); run.text = tt
                run.font.size = Pt(fs); run.font.bold = bd; run.font.name = font
                run.font.color.rgb = cc if cc else INK
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 else PANEL
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
    return h

def build(lang):
    L, F, font = G[lang], G[lang]['fig'], FONT[lang]
    prs = new_deck()

    # 표지
    sl = blank(prs)
    box(sl, 0, 0, SW, SH, fill=RGBColor(0x0F, 0x23, 0x3F))
    label(sl, 0, 2.3, SW, 0.9, 'PIE', size=64, color=RGBColor(0x60,0xA5,0xFA), bold=True,
          font=font, align=PP_ALIGN.CENTER)
    label(sl, 0, 3.25, SW, 0.6, L['h1'], size=26, color=WHITE, bold=True, font=font, align=PP_ALIGN.CENTER)
    label(sl, 1.5, 4.0, SW-3, 0.6, L['hsub'], size=14, color=RGBColor(0x93,0xC5,0xFD),
          font=font, align=PP_ALIGN.CENTER)
    label(sl, 0, 6.6, SW, 0.4, L['foot'], size=11, color=RGBColor(0x64,0x74,0x8B),
          font=font, align=PP_ALIGN.CENTER)

    # 목차
    sl = blank(prs); head(sl, L['tocT'], '', font)
    for i, s in enumerate(L['sec']):
        col, row = i // 9, i % 9
        label(sl, 0.8 + col*6.4, 1.35 + row*0.6, 6.0, 0.45,
              s['icon'] + '  ' + s['title'], size=14, font=font)

    # 각 절
    for si, s in enumerate(L['sec']):
        blocks, page = s['b'], 1
        sl = blank(prs)
        head(sl, s['icon'] + '  ' + s['title'], s.get('lead', ''), font, '%d/17' % (si+1))
        y = 1.25
        for blk in blocks:
            t = blk['t']
            if t == 'shot':
                fn = FIGS.get(blk['k'])
                if not fn:
                    continue
                fh = FIG_H.get(blk['k'], 3.2)
                if y + fh > 7.0:                      # 그림은 넉넉한 자리에만
                    page += 1; sl = blank(prs)
                    head(sl, s['icon'] + '  ' + s['title'], '', font, '%d/17' % (si+1)); y = 1.25
                fn(sl, F, font, y + 0.3)
                y += fh
                label(sl, 0.6, y, 12.1, 0.32, blk['cap'], size=11.5, color=DIM, font=font,
                      align=PP_ALIGN.CENTER)
                y += 0.5
                continue
            need = {'p': 0.55, 'h3': 0.45, 'tip': 0.6, 'warn': 0.6, 'note': 0.6}.get(t, 0)
            if t == 'steps': need = 0.34*len(blk['x']) + 0.2
            if t == 'tab':   need = 0.34 + 0.32*len(blk['rows']) + 0.25
            if y + need > 6.9:
                page += 1; sl = blank(prs)
                head(sl, s['icon'] + '  ' + s['title'], '', font, '%d/17' % (si+1)); y = 1.25
            if t == 'p':
                tb = label(sl, 0.7, y, 11.9, 0.5, blk['x'], size=13.5, font=font, anchor=MSO_ANCHOR.TOP)
                y += 0.62
            elif t == 'h3':
                label(sl, 0.7, y, 11.9, 0.4, blk['x'], size=15, color=BLUE, bold=True, font=font)
                y += 0.5
            elif t in ('tip', 'warn', 'note'):
                fill = {'tip': RGBColor(0xE8,0xF8,0xEE), 'warn': RGBColor(0xFE,0xE9,0xE9),
                        'note': RGBColor(0xE8,0xF1,0xFE)}[t]
                line = {'tip': GREEN, 'warn': RED, 'note': BLUE}[t]
                ic = {'tip': '💡 ', 'warn': '⚠ ', 'note': 'ℹ '}[t]
                bx = box(sl, 0.7, y, 11.9, 0.62, fill=fill, line=line, rad=True)
                label(sl, 0.9, y, 11.5, 0.62, ic + blk['x'], size=12.5, font=font)
                y += 0.78
            elif t == 'steps':
                for k, st in enumerate(blk['x']):
                    cir = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), Inches(y+0.03),
                                              Inches(0.3), Inches(0.3))
                    cir.fill.solid(); cir.fill.fore_color.rgb = BLUE
                    cir.line.fill.background(); cir.shadow.inherit = False
                    label(sl, 0.75, y+0.03, 0.3, 0.3, str(k+1), size=12, color=WHITE, bold=True,
                          font=font, align=PP_ALIGN.CENTER)
                    label(sl, 1.2, y, 11.4, 0.34, st, size=13, font=font)
                    y += 0.36
                y += 0.15
            elif t == 'tab':
                h = add_table(sl, blk['head'], blk['rows'], 0.7, y, 11.9, font)
                y += h + 0.25
    path = os.path.join(OUT, FNAME[lang])
    prs.save(path)
    return path, len(prs.slides._sldIdLst)

for lg in ('ko', 'zh', 'vi'):
    p, n = build(lg)
    print('%-4s %-34s %3d slides  %6.1f KB' % (lg, os.path.basename(p), n, os.path.getsize(p)/1024))
